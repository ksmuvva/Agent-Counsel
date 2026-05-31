"""Document I/O tools (Excel, Word, PowerPoint).

Real implementations backed by openpyxl / python-docx / python-pptx. Each
method raises :class:`MissingDependencyError` with an actionable message when
the optional library is not installed, rather than silently faking data.

Two surfaces are provided:

* :class:`DocumentTools` — a thin static-method helper for direct Python use.
* Per-format :class:`~tools.base.Tool` subclasses (e.g. ``ExcelReadTool``)
  that agents can invoke through the schema-based ``ToolRegistry``.
"""
from __future__ import annotations

from typing import Any, Dict, List

from .base import Tool, ToolError


class MissingDependencyError(RuntimeError):
    """Raised when an optional document library is not installed."""


# ---------------------------------------------------------------------------
# Static helper (kept for direct Python use and backwards compatibility).
# ---------------------------------------------------------------------------
class DocumentTools:
    """Read and write common office document formats."""

    # ---- Excel ---------------------------------------------------------
    @staticmethod
    def read_excel(path: str) -> Dict[str, List[List[Any]]]:
        try:
            import openpyxl
        except ImportError as exc:  # pragma: no cover
            raise MissingDependencyError(
                "openpyxl is required to read Excel files (pip install openpyxl)."
            ) from exc
        wb = openpyxl.load_workbook(path, data_only=True)
        return {
            sheet.title: [list(row) for row in sheet.iter_rows(values_only=True)]
            for sheet in wb.worksheets
        }

    @staticmethod
    def write_excel(path: str, rows: List[List[Any]], sheet_name: str = "Sheet1") -> str:
        try:
            import openpyxl
        except ImportError as exc:  # pragma: no cover
            raise MissingDependencyError(
                "openpyxl is required to write Excel files (pip install openpyxl)."
            ) from exc
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name
        for row in rows:
            ws.append(row)
        wb.save(path)
        return path

    # ---- Word ----------------------------------------------------------
    @staticmethod
    def read_word(path: str) -> str:
        try:
            import docx
        except ImportError as exc:  # pragma: no cover
            raise MissingDependencyError(
                "python-docx is required to read Word files (pip install python-docx)."
            ) from exc
        document = docx.Document(path)
        return "\n".join(p.text for p in document.paragraphs)

    @staticmethod
    def write_word(path: str, content: str) -> str:
        try:
            import docx
        except ImportError as exc:  # pragma: no cover
            raise MissingDependencyError(
                "python-docx is required to write Word files (pip install python-docx)."
            ) from exc
        document = docx.Document()
        for paragraph in content.split("\n"):
            document.add_paragraph(paragraph)
        document.save(path)
        return path

    # ---- PowerPoint ----------------------------------------------------
    @staticmethod
    def read_powerpoint(path: str) -> List[str]:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise MissingDependencyError(
                "python-pptx is required to read PowerPoint files (pip install python-pptx)."
            ) from exc
        prs = Presentation(path)
        slides: List[str] = []
        for slide in prs.slides:
            texts = [
                shape.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text
            ]
            slides.append("\n".join(texts))
        return slides

    @staticmethod
    def write_powerpoint(path: str, slides: List[Dict[str, str]]) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise MissingDependencyError(
                "python-pptx is required to write PowerPoint files (pip install python-pptx)."
            ) from exc
        prs = Presentation()
        layout = prs.slide_layouts[1]
        for spec in slides:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = spec.get("title", "")
            slide.placeholders[1].text = spec.get("body", "")
        prs.save(path)
        return path


# ---------------------------------------------------------------------------
# Schema-based tool wrappers.
# ---------------------------------------------------------------------------
class ExcelReadTool(Tool):
    name = "excel_read"
    description = (
        "Read every sheet of an Excel workbook (.xlsx) and return its cells "
        "as a dict mapping sheet name to a list of rows."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Filesystem path to the .xlsx file."}
        },
        "required": ["path"],
        "additionalProperties": False,
    }

    def execute(self, *, path: str) -> Dict[str, List[List[Any]]]:
        return DocumentTools.read_excel(path)


class ExcelWriteTool(Tool):
    name = "excel_write"
    description = "Write rows to a single-sheet Excel workbook (.xlsx)."
    input_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "rows": {
                "type": "array",
                "items": {"type": "array"},
                "description": "List of rows; each row is a list of cell values.",
            },
            "sheet_name": {"type": "string", "default": "Sheet1"},
        },
        "required": ["path", "rows"],
        "additionalProperties": False,
    }

    def execute(self, *, path: str, rows: List[List[Any]], sheet_name: str = "Sheet1") -> str:
        return DocumentTools.write_excel(path, rows, sheet_name)


class WordReadTool(Tool):
    name = "word_read"
    description = "Read a Word document (.docx) and return its body text."
    input_schema = {
        "type": "object",
        "properties": {"path": {"type": "string"}},
        "required": ["path"],
        "additionalProperties": False,
    }

    def execute(self, *, path: str) -> str:
        return DocumentTools.read_word(path)


class WordWriteTool(Tool):
    name = "word_write"
    description = "Write paragraphs to a Word document (.docx). Newlines become paragraph breaks."
    input_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "content": {"type": "string"},
        },
        "required": ["path", "content"],
        "additionalProperties": False,
    }

    def execute(self, *, path: str, content: str) -> str:
        return DocumentTools.write_word(path, content)


class PowerPointReadTool(Tool):
    name = "powerpoint_read"
    description = "Read a PowerPoint deck (.pptx) and return a list of slide texts."
    input_schema = {
        "type": "object",
        "properties": {"path": {"type": "string"}},
        "required": ["path"],
        "additionalProperties": False,
    }

    def execute(self, *, path: str) -> List[str]:
        return DocumentTools.read_powerpoint(path)


class PowerPointWriteTool(Tool):
    name = "powerpoint_write"
    description = "Write a PowerPoint deck (.pptx) from a list of {title, body} slides."
    input_schema = {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "body": {"type": "string"},
                    },
                    "additionalProperties": False,
                },
            },
        },
        "required": ["path", "slides"],
        "additionalProperties": False,
    }

    def execute(self, *, path: str, slides: List[Dict[str, str]]) -> str:
        return DocumentTools.write_powerpoint(path, slides)


# ---------------------------------------------------------------------------
# Web search.
# ---------------------------------------------------------------------------
class WebSearchTool(Tool):
    """Web research tool, backed by Tavily when ``TAVILY_API_KEY`` is set."""

    name = "web_search"
    description = (
        "Search the public web and return up to N {title, url, content} results. "
        "Requires TAVILY_API_KEY."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "query": {"type": "string", "minLength": 1},
            "max_results": {"type": "integer", "minimum": 1, "maximum": 20, "default": 5},
        },
        "required": ["query"],
        "additionalProperties": False,
    }

    def execute(self, *, query: str, max_results: int = 5) -> List[Dict[str, str]]:
        import os

        api_key = os.getenv("TAVILY_API_KEY")
        if not api_key:
            raise ToolError(
                "Web search requires TAVILY_API_KEY (or wire in your own provider)."
            )
        try:
            import requests
        except ImportError as exc:  # pragma: no cover
            raise ToolError(
                "requests is required for web search (pip install requests)."
            ) from exc
        resp = requests.post(
            "https://api.tavily.com/search",
            json={"api_key": api_key, "query": query, "max_results": max_results},
            timeout=30,
        )
        resp.raise_for_status()
        data = resp.json()
        return [
            {"title": r.get("title", ""), "url": r.get("url", ""), "content": r.get("content", "")}
            for r in data.get("results", [])
        ]
